from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from io import BytesIO
from typing import Optional


class SlideGenerator:
    def __init__(self):
        self.prs = Presentation()
        # Remove default empty slide
        xml_slides = self.prs.slides._sldIdLst
        if len(xml_slides) > 0:
            xml_slides.remove(xml_slides[0])
        self.slide_count = 0

    def add_title_slide(self, title: str, subtitle: str = "") -> bool:
        """Add title slide (Layout 0)"""
        try:
            layout = self.prs.slide_layouts[0]
            slide = self.prs.slides.add_slide(layout)

            slide.shapes.title.text = title
            if subtitle and len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle

            self.slide_count += 1
            return True
        except Exception as e:
            print(f"❌ Error adding title slide: {e}")
            return False

    def add_content_slide(self, heading: str, bullet_points: list, image_path: Optional[str] = None) -> bool:
        """Add content slide with bullets and optional image (Layout 1)"""
        try:
            layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(layout)

            # Set heading
            slide.shapes.title.text = heading

            # Add bullet points
            if bullet_points:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()

                for i, point in enumerate(bullet_points):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = point
                    p.font.size = Pt(18)
                    p.font.name = "Calibri"
                    p.space_after = Pt(8)
                    if i > 0:
                        p.level = 1  # Indent as bullet

            # Add image if path exists
            if image_path and os.path.exists(image_path):
                self._add_image_to_slide(slide, image_path)

            self.slide_count += 1
            return True
        except Exception as e:
            print(f"❌ Error adding content slide: {e}")
            return False

    def _add_image_to_slide(self, slide, image_path: str):
        """Place image on right side of slide"""
        try:
            left = Inches(5.5)
            top = Inches(1.5)
            width = Inches(4.0)
            height = Inches(5.0)
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            print(f"⚠️ Warning: Could not add image to slide: {e}")

    def save_to_file(self, filepath: str) -> str:
        """Save PPTX to disk"""
        if not filepath.lower().endswith(".pptx"):
            filepath += ".pptx"

        # Ensure directory exists
        os.makedirs(os.path.dirname(filepath) if os.path.dirname(filepath) else ".", exist_ok=True)
        self.prs.save(filepath)
        return filepath

    def save_to_bytes(self) -> BytesIO:
        """Save PPTX to BytesIO for Streamlit direct download"""
        bio = BytesIO()
        self.prs.save(bio)
        bio.seek(0)
        return bio

    def reset(self):
        """Clear all slides and restart"""
        self.prs = Presentation()
        xml_slides = self.prs.slides._sldIdLst
        if len(xml_slides) > 0:
            xml_slides.remove(xml_slides[0])
        self.slide_count = 0